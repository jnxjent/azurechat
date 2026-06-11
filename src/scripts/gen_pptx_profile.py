#!/usr/bin/env python3
"""
gen_pptx_profile.py
会社紹介PPT生成スクリプト（python-pptx + cairosvg アイコン付き）
Usage: python3 gen_pptx_profile.py --input data.json --output output.pptx [--palette navy_orange]
"""
import argparse, io, json, re, sys

try:
    import cairosvg
    HAS_CAIRO = True
except ImportError:
    HAS_CAIRO = False
    print("[gen_pptx_profile] cairosvg not available; icons skipped", file=sys.stderr)

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Slide geometry (pptxgenjs LAYOUT_WIDE と同じ 13.33×7.5 inch) ──────────
EMU = 914400
SW  = int(13.333 * EMU)
SH  = int(7.5   * EMU)
HH  = int(1.05  * EMU)   # header height

def inch(v): return int(v * EMU)
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ── カラーパレット定義 ─────────────────────────────────────────────────────
# 役割キー（全パレット共通）
#   main         : 帯・見出し・本文・番号円（濃色）
#   accent       : 強調アクセント（円・ストリップ・CTA）
#   accent_light : コールアウトボックス塗り（accentの淡色）
#   main_light   : カード枠・薄面（mainの淡色）
#   text_muted   : 補足・キャプションのグレー文字
# 背景・帯上の文字は全パレット共通で FFFFFF 固定。
PALETTES = {
    'navy_orange':    {'main': '13294B', 'accent': 'F5821F', 'accent_light': 'FBEDE0', 'main_light': 'E4E8F0', 'text_muted': '6B7488'},
    'forest_amber':   {'main': '1B4D3E', 'accent': 'F4A300', 'accent_light': 'FBEFD5', 'main_light': 'E3EDE8', 'text_muted': '5E6E66'},
    'burgundy_gold':  {'main': '8C1D18', 'accent': 'E0A33B', 'accent_light': 'F7ECD6', 'main_light': 'F3E5E4', 'text_muted': '6E5A58'},
    'teal_coral':     {'main': '0E4D5C', 'accent': 'EE6C4D', 'accent_light': 'FBE6DE', 'main_light': 'DCE9EC', 'text_muted': '5A6B70'},
    'charcoal_terra': {'main': '333333', 'accent': 'C15F3C', 'accent_light': 'F3E3DA', 'main_light': 'ECECEA', 'text_muted': '6E6E6E'},
}
DEFAULT_PALETTE = 'navy_orange'

FONT = 'Meiryo'

def _luminance(hex_color):
    """0–1 の相対輝度（RGB 加重平均）"""
    h = hex_color.lstrip('#')
    r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    return (0.299*r + 0.587*g + 0.114*b) / 255

def resolve_palette(key_or_dict):
    """
    パレット名(str) または役割キー辞書(dict) を受け取り、
    RGBColor 値を持つ描画用辞書 P を返す。
    将来のアドホック指定（main/accent直接渡し）にも対応。
    """
    src = key_or_dict if isinstance(key_or_dict, dict) \
          else PALETTES.get(key_or_dict, PALETTES[DEFAULT_PALETTE])

    # アクセント円内の文字色: accent が明るい(輝度>0.45)なら main、暗ければ白
    on_accent = rgb(src['main']) if _luminance(src['accent']) > 0.45 else rgb('FFFFFF')

    return {
        'main':         rgb(src['main']),
        'accent':       rgb(src['accent']),
        'accent_light': rgb(src['accent_light']),
        'main_light':   rgb(src['main_light']),
        'text_muted':   rgb(src['text_muted']),
        'white':        rgb('FFFFFF'),
        'on_accent':    on_accent,  # accent色の上に置く文字色（視認性確保）
    }

# 描画ヘルパーが参照するグローバル P（main() で resolve_palette() が設定する）
P: dict = {}

# ── SVG アイコン定義（Heroicons アウトラインスタイル、白ストローク）────────
_SVG_BASE = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" fill="none" stroke="white" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">{}</svg>'

ICONS = {
    "calendar":    _SVG_BASE.format('<rect x="3" y="4" width="18" height="18" rx="2"/><line x1="16" y1="2" x2="16" y2="6"/><line x1="8" y1="2" x2="8" y2="6"/><line x1="3" y1="10" x2="21" y2="10"/>'),
    "map-pin":     _SVG_BASE.format('<path d="M21 10c0 7-9 13-9 13s-9-6-9-13a9 9 0 0 1 18 0z"/><circle cx="12" cy="10" r="3"/>'),
    "trending-up": _SVG_BASE.format('<polyline points="23 6 13.5 15.5 8.5 10.5 2 17"/><polyline points="17 6 23 6 23 12"/>'),
    "users":       _SVG_BASE.format('<path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 0 0-3-3.87"/><path d="M16 3.13a4 4 0 0 1 0 7.75"/>'),
    "truck":       _SVG_BASE.format('<rect x="1" y="3" width="15" height="13"/><polygon points="16 8 20 8 23 11 23 16 16 16 16 8"/><circle cx="5.5" cy="18.5" r="2.5"/><circle cx="18.5" cy="18.5" r="2.5"/>'),
    "recycle":     _SVG_BASE.format('<polyline points="1.5 12 4.5 15 7.5 12"/><path d="M4.5 15V9a6 6 0 0 1 11.1-3.1"/><polyline points="22.5 12 19.5 9 16.5 12"/><path d="M19.5 9v6a6 6 0 0 1-11.1 3.1"/>'),
    "layers":      _SVG_BASE.format('<polygon points="12 2 2 7 12 12 22 7 12 2"/><polyline points="2 17 12 22 22 17"/><polyline points="2 12 12 17 22 12"/>'),
    "leaf":        _SVG_BASE.format('<path d="M2 22 16 8"/><path d="M16 8C14.5 5 12 2 8 2 5 2 2 5 2 8c0 4 3 6 6 8l2 2a12 12 0 0 0 6-10z"/>'),
    "award":       _SVG_BASE.format('<circle cx="12" cy="8" r="6"/><path d="M15.477 12.89 17 22l-5-3-5 3 1.523-9.11"/>'),
    "shield":      _SVG_BASE.format('<path d="M12 22s8-4 8-10V5l-8-3-8 3v7c0 6 8 10 8 10z"/>'),
    "check":       _SVG_BASE.format('<polyline points="20 6 9 17 4 12"/>'),
}

METRIC_ICON_MAP = {
    "創業": "calendar", "設立": "calendar",
    "本社": "map-pin",  "拠点": "map-pin",  "所在地": "map-pin",
    "上場": "trending-up", "株式": "trending-up",
    "取引": "users",    "顧客": "users",    "従業員": "users",
}
STEP_ICON_MAP = {
    "収集": "truck",  "運搬": "truck",
    "中間": "recycle", "処理": "recycle",
    "最終": "layers",  "処分": "layers",
}

def _icon_key(label, mapping):
    for k, v in mapping.items():
        if k in label:
            return v
    return "check"

def _svg_png(name, size=48):
    if not HAS_CAIRO:
        return None
    svg = ICONS.get(name, ICONS["check"])
    try:
        return cairosvg.svg2png(bytestring=svg.encode(), output_width=size, output_height=size)
    except Exception as e:
        print(f"[gen_pptx_profile] svg→png failed ({name}): {e}", file=sys.stderr)
        return None

# ── 描画ヘルパー ──────────────────────────────────────────────────────────
RECT = 1
ROUNDED_RECT = 5
OVAL = 9
RIGHT_ARROW = 13

def _shape(slide, shape_id, x, y, w, h, fill=None, line_color=None, line_pt=0.5, alpha=None):
    s = slide.shapes.add_shape(shape_id, int(x), int(y), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha is not None:
            s.fill.fore_color._element.set('lumMod', str(int(alpha * 100000)))
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def _text(slide, text, x, y, w, h, size, color, bold=False,
          align='left', valign='top', wrap=True, font=FONT):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    # East Asian フォントも同じ書体に設定（日本語が確実に Meiryo で描画されるよう）
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', {'top': 't', 'middle': 'ctr', 'bottom': 'b'}.get(valign, 't'))
    return tb

def _icon(slide, name, x, y, size_inch):
    size_px = max(int(size_inch * 96), 32)
    png = _svg_png(name, size_px)
    if png is None:
        return
    sz = inch(size_inch)
    slide.shapes.add_picture(io.BytesIO(png), int(x), int(y), sz, sz)

def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def _header(slide, title, kicker=None):
    """帯：main色塗り（紺等）＋白文字タイトル。kicker があればaccentで上部に表示。"""
    _shape(slide, RECT, 0, 0, SW, HH, fill=P['main'])
    _shape(slide, RECT, 0, HH, SW, inch(0.05), fill=P['accent'])  # 帯下のアクセントライン
    if kicker:
        _text(slide, kicker, inch(0.45), inch(0.08), SW - inch(0.9), inch(0.26),
              9, P['accent'], bold=True)
    ty = inch(0.3) if kicker else inch(0.15)
    _text(slide, title, inch(0.45), ty, SW - inch(0.9), HH - ty - inch(0.1),
          22, P['white'], bold=True)

def _chrome(slide):
    """スライド下端の細いアクセントライン"""
    _shape(slide, RECT, 0, SH - inch(0.08), SW, inch(0.08), fill=P['accent'])

# ── スライドビルダー ───────────────────────────────────────────────────────

def build_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, P['main'])  # 表紙は main 色の濃色背景

    # 装飾サークル（右上）— accent + 淡い変化でレイヤー感
    _shape(slide, OVAL, SW - inch(3.5), -inch(0.6), inch(3.8), inch(3.8), fill=P['accent'])
    _shape(slide, OVAL, SW - inch(2.6),  inch(1.1),  inch(2.8), inch(2.8), fill=P['accent'])
    _shape(slide, OVAL, SW - inch(3.8),  inch(1.6),  inch(1.8), inch(1.8), fill=P['main_light'])

    # 左縦バー + 下部バー
    _shape(slide, RECT, 0, 0, inch(0.28), SH, fill=P['accent'])
    _shape(slide, RECT, 0, SH - inch(0.12), SW, inch(0.12), fill=P['accent'])

    # Kicker + 区切り線
    _text(slide, 'COMPANY PROFILE', inch(0.72), inch(0.85), inch(5.5), inch(0.28),
          10, P['accent'], bold=True)
    _shape(slide, RECT, inch(0.72), inch(1.18), inch(0.72), inch(0.05), fill=P['accent'])

    # タイトル（左揃え）
    title = data.get('title', '会社紹介')
    _text(slide, title, inch(0.72), inch(1.35), inch(8.0), inch(2.2),
          32, P['white'], bold=True, valign='middle')

    # サブタイトル
    subtitle = data.get('coverSubtitle', '')
    if subtitle:
        _text(slide, subtitle, inch(0.72), inch(3.75), inch(8.0), inch(0.55),
              14, P['accent'])

    # ページ数
    n = data.get('slideCount', 7)
    _text(slide, f'{n} pages', SW - inch(2.2), SH - inch(0.45), inch(1.8), inch(0.25),
          10, P['white'], bold=True, align='right')


def build_company_overview(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, P['white'])  # 本文スライドは白背景
    _header(slide, sd.get('title', '会社概要'), 'ABOUT US')

    cy = HH + inch(0.22)
    ch = SH - cy - inch(0.22)
    lw = inch(5.1)
    rx = inch(5.72)
    rw = SW - rx - inch(0.42)

    # ── 左パネル ──
    lead = sd.get('leadText', '')
    if lead:
        _text(slide, lead, inch(0.48), cy, lw, inch(2.6), 12, P['main'])

    callout = sd.get('callout')
    if callout:
        coy = cy + inch(2.75)
        coh = max(ch - inch(2.75) - inch(0.1), inch(1.8))
        _shape(slide, ROUNDED_RECT, inch(0.48), coy, lw, coh,
               fill=P['accent_light'], line_color=P['accent'], line_pt=1.2)
        _shape(slide, RECT, inch(0.48), coy, inch(0.22), coh, fill=P['accent'])
        _text(slide, callout['title'], inch(0.82), coy + inch(0.14),
              lw - inch(0.48), inch(0.32), 11, P['accent'], bold=True)
        _text(slide, callout['body'], inch(0.82), coy + inch(0.5),
              lw - inch(0.48), coh - inch(0.68), 10, P['main'])

    # 縦区切り線
    _shape(slide, RECT, inch(5.62), cy + inch(0.1), inch(0.04), ch - inch(0.2), fill=P['main_light'])

    # ── 右パネル: 2×2 メトリクスカード ──
    metrics = sd.get('metrics', [])
    gap = inch(0.14)
    cw  = (rw - gap) / 2
    cah = (ch - gap) / 2

    for i, m in enumerate(metrics[:4]):
        col = i % 2
        row = i // 2
        x = rx + col * (cw + gap)
        y = cy + row * (cah + gap)

        # カード背景
        _shape(slide, ROUNDED_RECT, x, y, cw, cah,
               fill=P['white'], line_color=P['main_light'], line_pt=0.8)
        # 左アクセントストリップ
        _shape(slide, ROUNDED_RECT, x, y, inch(0.22), cah, fill=P['accent'])

        # アイコン（アクセントストリップの上部）
        icon_name = _icon_key(m.get('label', ''), METRIC_ICON_MAP)
        _icon(slide, icon_name, x + inch(0.02), y + inch(0.1), 0.18)

        lbl_x = x + inch(0.32)
        lbl_w = cw - inch(0.42)
        _text(slide, m.get('label', ''),
              lbl_x, y + inch(0.15), lbl_w, inch(0.28), 10, P['text_muted'])
        _text(slide, m.get('value', ''),
              lbl_x, y + inch(0.42), lbl_w, cah * 0.42, 22, P['main'],
              bold=True, valign='middle')
        if m.get('note'):
            _text(slide, m['note'], lbl_x, y + cah - inch(0.38), lbl_w, inch(0.28),
                  9, P['text_muted'])

    _chrome(slide)


def build_process_cards(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, P['white'])
    _header(slide, sd.get('title', '一貫体制'), 'OUR STRENGTH')

    cy = HH + inch(0.18)
    sub = sd.get('subtitle', '')
    if sub:
        _text(slide, sub, inch(0.48), cy, SW - inch(0.96), inch(0.4), 11, P['text_muted'])
        cy += inch(0.44)

    steps    = sd.get('steps', [])
    benefits = sd.get('benefits', [])
    ben_h    = inch(0.52) if benefits else 0
    cards_h  = SH - cy - ben_h - inch(0.26)
    count    = min(len(steps), 4)
    gap      = inch(0.14)
    tw       = SW - inch(0.84)
    cw       = (tw - gap * (count - 1)) / count
    icon_d   = inch(0.68)

    for i, step in enumerate(steps[:count]):
        cx = inch(0.42) + i * (cw + gap)
        # 最初のカードは accent、以降は main で視覚的階層を表現
        af         = P['accent'] if i == 0 else P['main']
        text_on_af = P['on_accent'] if i == 0 else P['white']
        num = f'{i+1:02d}'

        # カード
        _shape(slide, ROUNDED_RECT, cx, cy, cw, cards_h,
               fill=P['main_light'], line_color=P['main_light'], line_pt=0.8)

        # ウォーターマーク数字（右上、薄め）
        _text(slide, num, cx + cw - inch(1.0), cy + inch(0.04),
              inch(0.9), inch(0.82), 36, P['text_muted'], bold=True, align='right')

        # アイコン円（中央上部）
        ix = cx + cw / 2 - icon_d / 2
        iy = cy + inch(0.2)
        _shape(slide, OVAL, ix, iy, icon_d, icon_d, fill=af)

        icon_name = _icon_key(step.get('title', ''), STEP_ICON_MAP)
        _icon(slide, icon_name, ix + inch(0.1), iy + inch(0.1), 0.48)

        # タイトル
        ty = iy + icon_d + inch(0.14)
        _text(slide, step.get('title', ''), cx + inch(0.1), ty,
              cw - inch(0.2), inch(0.44), 14, P['main'],
              bold=True, align='center', valign='middle')

        # 本文
        by = ty + inch(0.48)
        _text(slide, step.get('body', ''), cx + inch(0.14), by,
              cw - inch(0.28), cards_h - (by - cy) - inch(0.16),
              11, P['text_muted'])

        # 矢印コネクタ
        if i < count - 1:
            ax = cx + cw + gap * 0.1
            ah = cards_h * 0.13
            ay = cy + cards_h / 2 - ah / 2
            _shape(slide, RIGHT_ARROW, ax, ay, gap * 0.68, ah, fill=P['accent'])

    # ── 下部メリット行 ──
    if benefits:
        bar_y = SH - ben_h - inch(0.18)
        _shape(slide, RECT, inch(0.42), bar_y - inch(0.04),
               SW - inch(0.84), inch(0.04), fill=P['accent'])
        iw = (SW - inch(0.84)) / (len(benefits) + 1)
        _text(slide, '一貫体制のメリット：', inch(0.42), bar_y,
              iw, ben_h - inch(0.06), 10, P['accent'], bold=True, valign='middle')
        for j, b in enumerate(benefits):
            _text(slide, f'✓ {b}', inch(0.42) + iw * (j + 1), bar_y,
                  iw, ben_h - inch(0.06), 10, P['main'], valign='middle')

    _chrome(slide)


def _split_bullet(text):
    m = re.match(r'^([^：:]{2,18})[：:]\s*([\s\S]*)$', text)
    if m:
        return m.group(1), m.group(2)
    return text[:14], text[14:]


def build_bullets(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, P['white'])
    _header(slide, sd.get('title', ''))

    bullets = sd.get('bullets', [])
    count = min(len(bullets), 5)
    if count == 0:
        _chrome(slide)
        return

    cy   = HH + inch(0.18)
    totH = SH - cy - inch(0.26)

    if count <= 3:
        # 横並びカード
        gap  = inch(0.14)
        tw   = SW - inch(0.84)
        cw   = (tw - gap * (count - 1)) / count
        icD  = inch(0.68)

        for i, bullet in enumerate(bullets[:count]):
            cx = inch(0.42) + i * (cw + gap)
            af         = P['accent'] if i == 0 else P['main']
            text_on_af = P['on_accent'] if i == 0 else P['white']
            num = f'{i+1:02d}'
            t, body = _split_bullet(bullet)

            _shape(slide, ROUNDED_RECT, cx, cy, cw, totH,
                   fill=P['main_light'], line_color=P['main_light'], line_pt=0.8)
            _text(slide, num, cx + cw - inch(1.0), cy + inch(0.04),
                  inch(0.9), inch(0.82), 36, P['text_muted'], bold=True, align='right')

            # アイコン円
            ix = cx + cw / 2 - icD / 2
            iy = cy + inch(0.2)
            _shape(slide, OVAL, ix, iy, icD, icD, fill=af)
            _text(slide, num, ix + inch(0.04), iy + inch(0.04),
                  icD - inch(0.08), icD - inch(0.08), 14, text_on_af,
                  bold=True, align='center', valign='middle')

            ty = iy + icD + inch(0.14)
            _text(slide, t, cx + inch(0.1), ty, cw - inch(0.2), inch(0.44),
                  13, P['main'], bold=True, align='center', valign='middle')
            _text(slide, body, cx + inch(0.14), ty + inch(0.5),
                  cw - inch(0.28), totH - (ty - cy) - inch(0.58),
                  11, P['text_muted'])

            if i < count - 1:
                ah = totH * 0.13
                _shape(slide, RIGHT_ARROW,
                       cx + cw + gap * 0.1, cy + totH / 2 - ah / 2,
                       gap * 0.68, ah, fill=P['accent'])
    else:
        # 縦並びカード
        gap  = inch(0.1)
        cah  = max(inch(0.88), (totH - gap * (count - 1)) / count)
        sw   = inch(0.38)
        totW = SW - inch(0.84)

        for i, bullet in enumerate(bullets[:count]):
            y  = cy + i * (cah + gap)
            af         = P['accent'] if i == 0 else P['main']
            text_on_af = P['on_accent'] if i == 0 else P['white']
            num = f'{i+1:02d}'

            _shape(slide, ROUNDED_RECT, inch(0.42), y, totW, cah,
                   fill=P['white'], line_color=P['main_light'], line_pt=0.8)
            _shape(slide, ROUNDED_RECT, inch(0.42), y, sw, cah, fill=af)
            _text(slide, num, inch(0.44), y + cah / 2 - inch(0.22),
                  sw - inch(0.06), inch(0.44), 13, text_on_af,
                  bold=True, align='center', valign='middle')
            _text(slide, bullet, inch(0.42) + sw + inch(0.16), y + inch(0.1),
                  totW - sw - inch(0.24), cah - inch(0.2),
                  13, P['main'], valign='middle')

    _chrome(slide)


def build_closing(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, P['main'])  # クロージングは main 色の濃色背景

    # 装飾サークル（右上）
    _shape(slide, OVAL, SW - inch(3.5), -inch(0.6), inch(3.8), inch(3.8), fill=P['accent'])
    _shape(slide, OVAL, SW - inch(2.6),  inch(1.1),  inch(2.8), inch(2.8), fill=P['accent'])

    _shape(slide, RECT, 0, 0, inch(0.28), SH, fill=P['accent'])
    _shape(slide, RECT, 0, SH - inch(0.12), SW, inch(0.12), fill=P['accent'])

    _text(slide, sd.get('title', 'お問い合わせ'),
          inch(0.72), inch(1.8), SW - inch(4.0), inch(1.4),
          28, P['white'], bold=True, valign='middle')

    for i, b in enumerate(sd.get('bullets', [])[:4]):
        iy = inch(3.4) + i * inch(0.78)
        # accent_light 塗り + accent 枠：暗い背景に映えるハイライトカード
        _shape(slide, ROUNDED_RECT, inch(0.72), iy, SW - inch(4.5), inch(0.66),
               fill=P['accent_light'], line_color=P['accent'], line_pt=0.8)
        _text(slide, b, inch(0.94), iy + inch(0.08), SW - inch(4.9), inch(0.5),
              13, P['main'], valign='middle')


# ── メイン ────────────────────────────────────────────────────────────────
BUILDERS = {
    'company-overview': build_company_overview,
    'process-cards':    build_process_cards,
    'closing':          build_closing,
}

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--input',   required=True)
    ap.add_argument('--output',  required=True)
    ap.add_argument('--palette', default=DEFAULT_PALETTE,
                    choices=list(PALETTES.keys()),
                    help='カラーパレット名 (既定: %(default)s)')
    args = ap.parse_args()

    # グローバル P をパレット選択で初期化
    global P
    P = resolve_palette(args.palette)
    print(f"[gen_pptx_profile] palette={args.palette}", file=sys.stderr)

    with open(args.input, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)

    build_cover(prs, data)

    for sd in data.get('slides', []):
        lt = sd.get('layoutType', 'bullets')
        builder = BUILDERS.get(lt, build_bullets)
        builder(prs, sd)

    prs.save(args.output)
    print(json.dumps({"ok": True, "slides": len(prs.slides), "palette": args.palette}))

if __name__ == '__main__':
    main()
