"""
PPTX → PNG 変換スクリプト
- Windows: PowerPoint COM経由でPDF変換 → fitz でPNG化
- Linux:   LibreOffice headless でPDF変換 → fitz でPNG化
          ※ LibreOffice が LAYOUT_WIDE(13.33"×7.5") を正しく扱わない問題に対応:
            python-pptx で 10"×7.5" にスケールダウンした一時 PPTX を変換し
            レンダリング後に元のアスペクト比に戻す。
使用方法: python3 pptx_to_png.py <input.pptx> <output_dir> [max_slides]
"""
import sys
import os
import subprocess
import tempfile
import shutil
import platform

# LAYOUT_WIDE のスライドサイズ (PptxGenJS の pptx.layout = "LAYOUT_WIDE")
WIDE_W_IN = 13.333  # inches
WIDE_H_IN = 7.5     # inches
STD_W_IN  = 10.0    # inches (LibreOffice が安全に扱える標準幅)
STD_H_IN  = 7.5     # inches


def get_pptx_slide_dims_inches(pptx_path: str) -> tuple:
    """PPTXのスライドサイズをインチで返す。失敗時はLAYOUT_WIDEのデフォルト値。"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        w_in = prs.slide_width / 914400
        h_in = prs.slide_height / 914400
        print(f"[pptx_to_png] slide dims: {w_in:.3f}\" × {h_in:.3f}\"", file=sys.stderr)
        return w_in, h_in
    except Exception as e:
        print(f"[pptx_to_png] Could not read slide dims: {e}", file=sys.stderr)
        return WIDE_W_IN, WIDE_H_IN


def create_scaled_pptx(src_path: str, dst_path: str,
                       target_w_in: float, target_h_in: float) -> bool:
    """
    python-pptxで全シェイプ座標をtargetサイズに比例スケールした
    一時コピーを作成する。LibreOfficeへ渡す前処理として使用。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(src_path)
        orig_w = prs.slide_width   # EMU
        orig_h = prs.slide_height  # EMU
        tgt_w  = Inches(target_w_in)
        tgt_h  = Inches(target_h_in)

        scale_x = tgt_w / orig_w
        scale_y = tgt_h / orig_h

        if abs(scale_x - 1.0) < 0.01 and abs(scale_y - 1.0) < 0.01:
            print("[pptx_to_png] slide is already standard size, no scaling needed",
                  file=sys.stderr)
            return False  # スケール不要

        print(f"[pptx_to_png] scaling PPTX: {orig_w/914400:.3f}\"×{orig_h/914400:.3f}\" "
              f"→ {target_w_in:.3f}\"×{target_h_in:.3f}\" "
              f"(sx={scale_x:.3f}, sy={scale_y:.3f})",
              file=sys.stderr)

        prs.slide_width  = tgt_w
        prs.slide_height = tgt_h

        for slide in prs.slides:
            for shape in slide.shapes:
                shape.left   = int(shape.left   * scale_x)
                shape.top    = int(shape.top    * scale_y)
                shape.width  = int(shape.width  * scale_x)
                shape.height = int(shape.height * scale_y)

        prs.save(dst_path)
        return True
    except Exception as e:
        print(f"[pptx_to_png] create_scaled_pptx error: {e}", file=sys.stderr)
        return False


def convert_pptx_to_pdf_windows(pptx_path: str, pdf_path: str) -> bool:
    try:
        import comtypes.client
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        try:
            abs_pptx = os.path.abspath(pptx_path)
            abs_pdf  = os.path.abspath(pdf_path)
            prs = powerpoint.Presentations.Open(abs_pptx, ReadOnly=True, Untitled=False, WithWindow=False)
            prs.SaveAs(abs_pdf, 32)  # 32 = ppSaveAsPDF
            prs.Close()
            return True
        finally:
            powerpoint.Quit()
    except Exception as e:
        print(f"[pptx_to_png] COM error: {e}", file=sys.stderr)
        return False


def convert_pptx_to_pdf_libreoffice(pptx_path: str, out_dir: str) -> str | None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("[pptx_to_png] LibreOffice not found", file=sys.stderr)
        return None

    # プロセスごとに独立したLOプロファイルを使用（並行実行時の競合回避）
    lo_profile = f"/tmp/lo_profile_{os.getpid()}"
    os.makedirs(lo_profile, exist_ok=True)

    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--norestore",
            "--nologo",
            f"-env:UserInstallation=file://{lo_profile}",
            "--convert-to", "pdf",
            "--outdir", out_dir,
            pptx_path,
        ],
        capture_output=True, text=True, timeout=120
    )

    shutil.rmtree(lo_profile, ignore_errors=True)

    if result.returncode != 0:
        print(f"[pptx_to_png] LibreOffice error: {result.stderr}", file=sys.stderr)
        return None
    pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(out_dir, pdf_name)
    return pdf_path if os.path.exists(pdf_path) else None


def pdf_to_pngs(pdf_path: str, output_dir: str, max_slides: int,
                pptx_w_in: float = WIDE_W_IN, pptx_h_in: float = WIDE_H_IN,
                render_w_in: float = WIDE_W_IN, render_h_in: float = WIDE_H_IN) -> list:
    """
    PDFをPNGに変換する。
    pptx_w_in/h_in: 元PPTXのスライドサイズ（出力PNG比率の基準）
    render_w_in/h_in: 実際に変換したPPTXのサイズ（LibreOfficeに渡したもの）
    """
    try:
        import fitz
        doc = fitz.open(pdf_path)
        png_paths = []

        expected_render_ar = render_w_in / render_h_in  # LibreOfficeに渡したサイズのAR

        for i in range(min(len(doc), max_slides)):
            page = doc[i]
            mb = page.mediabox
            actual_ar = mb.width / mb.height if mb.height > 0 else expected_render_ar

            ar_diff = abs(actual_ar - expected_render_ar) / expected_render_ar
            if ar_diff > 0.05:
                # それでもARがずれている場合はX方向スケールで補正
                scale_x = (expected_render_ar / actual_ar) * 2.0
                scale_y = 2.0
                print(
                    f"[pptx_to_png] slide {i}: residual AR mismatch "
                    f"(expected {expected_render_ar:.3f}, got {actual_ar:.3f}). "
                    f"Applying scale_x={scale_x:.3f}",
                    file=sys.stderr
                )
            else:
                scale_x = 2.0
                scale_y = 2.0

            mat = fitz.Matrix(scale_x, scale_y)
            pix = page.get_pixmap(matrix=mat)
            out_path = os.path.join(output_dir, f"slide_{i}.png")
            pix.save(out_path)
            png_paths.append(out_path)

        doc.close()
        print(f"[pptx_to_png] Generated {len(png_paths)} PNGs", file=sys.stderr)
        return png_paths
    except Exception as e:
        print(f"[pptx_to_png] fitz error: {e}", file=sys.stderr)
        return []


def convert_pptx_to_pngs(pptx_path: str, output_dir: str, max_slides: int = 8) -> list:
    os.makedirs(output_dir, exist_ok=True)
    tmp_dir = tempfile.mkdtemp()
    try:
        pptx_w_in, pptx_h_in = get_pptx_slide_dims_inches(pptx_path)
        needs_scale = pptx_w_in > STD_W_IN + 0.1  # LAYOUT_WIDE など標準幅超え

        if platform.system() == "Windows":
            print("[pptx_to_png] Using PowerPoint COM (Windows)", file=sys.stderr)
            pdf_path = os.path.join(tmp_dir, "output.pdf")
            ok = convert_pptx_to_pdf_windows(pptx_path, pdf_path)
            if not ok:
                return []
            return pdf_to_pngs(pdf_path, output_dir, max_slides,
                               pptx_w_in, pptx_h_in, pptx_w_in, pptx_h_in)

        else:
            print("[pptx_to_png] Using LibreOffice (Linux)", file=sys.stderr)

            # LAYOUT_WIDE (13.33") → LibreOfficeが誤って10"として扱う問題を回避:
            # python-pptx で 10"×7.5" スケールダウンコピーを作成してから変換する
            if needs_scale:
                scaled_pptx = os.path.join(tmp_dir, "scaled.pptx")
                ok = create_scaled_pptx(pptx_path, scaled_pptx, STD_W_IN, STD_H_IN)
                lo_input = scaled_pptx if ok else pptx_path
                render_w = STD_W_IN if ok else pptx_w_in
                render_h = STD_H_IN if ok else pptx_h_in
            else:
                lo_input  = pptx_path
                render_w  = pptx_w_in
                render_h  = pptx_h_in

            result = convert_pptx_to_pdf_libreoffice(lo_input, tmp_dir)
            if not result:
                return []

            return pdf_to_pngs(result, output_dir, max_slides,
                               pptx_w_in, pptx_h_in, render_w, render_h)

    except Exception as e:
        print(f"[pptx_to_png] Exception: {e}", file=sys.stderr)
        return []
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: pptx_to_png.py <input.pptx> <output_dir> [max_slides]")
        sys.exit(1)

    pptx_path  = sys.argv[1]
    output_dir = sys.argv[2]
    max_slides = int(sys.argv[3]) if len(sys.argv) > 3 else 8

    paths = convert_pptx_to_pngs(pptx_path, output_dir, max_slides)
    for p in paths:
        print(p)
